#!/usr/bin/env python3
"""FlyAI - Professionelle PowerPoint Praesentation (v2 - mit allen Updates)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Farben (Liquid Glass Palette) ──
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x12, 0x12, 0x1C)
ACCENT_BLUE = RGBColor(0x5A, 0xC8, 0xFA)
ACCENT_PURPLE = RGBColor(0xBF, 0x5A, 0xF2)
ACCENT_PINK = RGBColor(0xFF, 0x37, 0x5F)
ACCENT_GREEN = RGBColor(0x30, 0xD1, 0x58)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
GLASS_BG = RGBColor(0x1A, 0x1A, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color=GLASS_BG, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=TEXT_DIM):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc, accent=ACCENT_BLUE):
    card = add_shape(slide, left, top, width, height, GLASS_BG, 0.08)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    add_text(slide, icon, left + Inches(0.2), top + Inches(0.4), Inches(1), Inches(0.6),
             font_size=28, color=accent, bold=True)
    add_text(slide, title, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.4),
             font_size=16, color=WHITE, bold=True)
    add_text(slide, desc, left + Inches(0.2), top + Inches(1.35), width - Inches(0.4), Inches(0.8),
             font_size=12, color=TEXT_DIM)
    return card


def slide_header(slide, label, title, label_color=ACCENT_BLUE):
    add_text(slide, label, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             font_size=14, color=label_color, bold=True)
    add_text(slide, title, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
             font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(1.5), label_color)


# ═══════════════════════════════════════════════════
# SLIDE 1: Titelfolie
# ═══════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1)

circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-1), Inches(6), Inches(6))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x15, 0x1A, 0x30)
circle1.line.fill.background()

circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(3), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x1A, 0x10, 0x2A)
circle2.line.fill.background()

badge = add_shape(slide1, Inches(5.4), Inches(1.2), Inches(2.5), Inches(0.55), RGBColor(0x15, 0x20, 0x35), 0.3)
add_text(slide1, "KI-gestuetzte Reiseplattform", Inches(5.4), Inches(1.22),
         Inches(2.5), Inches(0.55), font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text(slide1, "FlyAI", Inches(1), Inches(2.2), Inches(11.3), Inches(1.4),
         font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide1, "Intelligente Reiseplattform mit KI-Chatbot,\nFlug- & Hotelsuche, Visa-Check & Reiseberatung",
         Inches(2.5), Inches(3.5), Inches(8.3), Inches(1),
         font_size=22, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_accent_line(slide1, Inches(5.6), Inches(4.6), Inches(2.1), ACCENT_BLUE)

techs = ["Python", "FastAPI", "OpenAI GPT-4", "Structured Output", "Liquid Glass UI"]
badge_x = 2.8
for tech in techs:
    w = len(tech) * 0.12 + 0.4
    add_shape(slide1, Inches(badge_x), Inches(5.0), Inches(w), Inches(0.45), RGBColor(0x15, 0x1A, 0x28), 0.3)
    add_text(slide1, tech, Inches(badge_x), Inches(5.02), Inches(w), Inches(0.45),
             font_size=11, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    badge_x += w + 0.15

add_text(slide1, "Aryan Sakhi  |  Masterschool GenAI  |  2026", Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
         font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 2: Problem & Loesung
# ═══════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)

add_text(slide2, "Das Problem", Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         font_size=14, color=ACCENT_PINK, bold=True)
add_text(slide2, "Warum FlyAI?", Inches(0.8), Inches(0.9), Inches(5), Inches(0.7),
         font_size=36, color=WHITE, bold=True)
add_accent_line(slide2, Inches(0.8), Inches(1.55), Inches(1.5), ACCENT_PINK)

problems = [
    ("Fragmentierte Suche", "Nutzer muessen zwischen 5+ Plattformen wechseln fuer Fluege, Hotels, Visa-Infos und Reisewarnungen."),
    ("Keine intelligente Beratung", "Keine Plattform kombiniert KI-Chat mit Echtzeit-Flugsuche und automatischen Visa-Checks."),
    ("Info-Overload", "Reiseplanung ist zeitaufwendig: Budget, Packliste, beste Reisezeit — alles verstreut."),
]

y = 2.1
for title, desc in problems:
    card = add_shape(slide2, Inches(0.8), Inches(y), Inches(5.5), Inches(1.35), GLASS_BG, 0.06)
    ind = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Pt(4), Inches(1.35))
    ind.fill.solid()
    ind.fill.fore_color.rgb = ACCENT_PINK
    ind.line.fill.background()
    add_text(slide2, title, Inches(1.2), Inches(y + 0.15), Inches(4.8), Inches(0.35),
             font_size=16, color=WHITE, bold=True)
    add_text(slide2, desc, Inches(1.2), Inches(y + 0.55), Inches(4.8), Inches(0.7),
             font_size=12, color=TEXT_DIM)
    y += 1.5

add_text(slide2, "Die Loesung", Inches(7), Inches(0.5), Inches(5), Inches(0.6),
         font_size=14, color=ACCENT_GREEN, bold=True)
add_text(slide2, "Alles in einer App", Inches(7), Inches(0.9), Inches(5.5), Inches(0.7),
         font_size=36, color=WHITE, bold=True)
add_accent_line(slide2, Inches(7), Inches(1.55), Inches(1.5), ACCENT_GREEN)

solutions = [
    ("Zentrale Plattform", "Fluege, Hotels, Mietwagen und Aktivitaeten — alles an einem Ort suchen und vergleichen."),
    ("KI-Reiseberater", "GPT-4o-mini Chatbot mit 10+ Intent-Typen: Fluege, Hotels, Pakete, Budget, Packliste, Vergleiche."),
    ("Automatische Reise-Infos", "Visa-Check, Reisewarnungen vom Auswaertigen Amt, Sicherheitshinweise — alles automatisch."),
]

y = 2.1
for title, desc in solutions:
    card = add_shape(slide2, Inches(7), Inches(y), Inches(5.5), Inches(1.35), GLASS_BG, 0.06)
    ind = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(y), Pt(4), Inches(1.35))
    ind.fill.solid()
    ind.fill.fore_color.rgb = ACCENT_GREEN
    ind.line.fill.background()
    add_text(slide2, title, Inches(7.4), Inches(y + 0.15), Inches(4.8), Inches(0.35),
             font_size=16, color=WHITE, bold=True)
    add_text(slide2, desc, Inches(7.4), Inches(y + 0.55), Inches(4.8), Inches(0.7),
             font_size=12, color=TEXT_DIM)
    y += 1.5


# ═══════════════════════════════════════════════════
# SLIDE 3: Features (erweitert)
# ═══════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)
slide_header(slide3, "Features", "Was FlyAI kann")

features = [
    ("Flugsuche", "Echtzeit-Suche via Skyscanner API mit Filtern nach Airline, Preis, Stops und Reisedauer.", ACCENT_BLUE),
    ("Hotelsuche", "Hotels weltweit vergleichen mit Preisen, Bewertungen und Standort-Infos.", ACCENT_PURPLE),
    ("KI-Chatbot", "GPT-4o-mini mit 10+ Intents: Fluege, Hotels, Pakete, Inspiration, Vergleich, Packliste, Budget.", ACCENT_GREEN),
    ("Visa & Sicherheit", "Automatische Visa-Infos + Reisewarnungen vom Auswaertigen Amt fuer 45+ Laender.", ACCENT_ORANGE),
    ("Mietwagen", "Mietwagen-Suche mit Vergleich von Sixt, Europcar, Hertz und Enterprise.", ACCENT_PINK),
    ("Aktivitaeten", "Kultur, Abenteuer, Food-Tours — Aktivitaeten fuer jedes Reiseziel entdecken.", ACCENT_BLUE),
]

col1_x, col2_x, col3_x = 0.8, 4.95, 9.1
row1_y, row2_y = 2.1, 4.65
positions = [
    (col1_x, row1_y), (col2_x, row1_y), (col3_x, row1_y),
    (col1_x, row2_y), (col2_x, row2_y), (col3_x, row2_y),
]

for i, (title, desc, color) in enumerate(features):
    x, y = positions[i]
    add_icon_card(slide3, Inches(x), Inches(y), Inches(3.85), Inches(2.2),
                  str(i + 1), title, desc, color)


# ═══════════════════════════════════════════════════
# SLIDE 4: Architektur
# ═══════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)
slide_header(slide4, "Architektur", "System-Design")

# Frontend
fe_box = add_shape(slide4, Inches(0.8), Inches(2.2), Inches(3.5), Inches(4.5), RGBColor(0x0F, 0x17, 0x2A), 0.04)
add_text(slide4, "FRONTEND", Inches(0.8), Inches(2.3), Inches(3.5), Inches(0.4),
         font_size=12, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
fe_items = [
    "HTML5 / CSS3 / Vanilla JS",
    "Apple Liquid Glass Design",
    "Clash Display + Satoshi Fonts",
    "4 Seiten: Fluege, Hotels,",
    "  Mietwagen, Aktivitaeten",
    "Responsive (Mobile/Tablet)",
    "XSS-Schutz (textContent)",
    "Accessibility (ARIA Labels)",
]
add_bullet_list(slide4, fe_items, Inches(1.1), Inches(2.8), Inches(3), Inches(3.5), font_size=12)

arrow1 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(4.1), Inches(1), Inches(0.5))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = ACCENT_BLUE
arrow1.line.fill.background()

# Backend
be_box = add_shape(slide4, Inches(5.7), Inches(2.2), Inches(3.5), Inches(4.5), RGBColor(0x0F, 0x17, 0x2A), 0.04)
add_text(slide4, "BACKEND", Inches(5.7), Inches(2.3), Inches(3.5), Inches(0.4),
         font_size=12, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
be_items = [
    "Python 3.13 + FastAPI",
    "SQLModel + SQLite (4 Tabellen)",
    "20+ REST Endpunkte",
    "Rate Limiting (SlowAPI)",
    "Security Headers Middleware",
    "Input Sanitization (bleach)",
    "bcrypt Passwort-Hashing",
    "Structured Output (JSON Mode)",
]
add_bullet_list(slide4, be_items, Inches(6), Inches(2.8), Inches(3), Inches(3.5), font_size=12)

arrow2 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.4), Inches(4.1), Inches(1), Inches(0.5))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = ACCENT_PURPLE
arrow2.line.fill.background()

# APIs
api_box = add_shape(slide4, Inches(10.6), Inches(2.2), Inches(2.2), Inches(4.5), RGBColor(0x0F, 0x17, 0x2A), 0.04)
add_text(slide4, "EXTERNE APIs", Inches(10.6), Inches(2.3), Inches(2.2), Inches(0.4),
         font_size=12, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
api_items = [
    "OpenAI GPT-4o-mini",
    "Skyscanner API",
    "Sky-Scrapper API",
    "Auswaertiges Amt",
]
add_bullet_list(slide4, api_items, Inches(10.8), Inches(2.8), Inches(1.8), Inches(3.5), font_size=12)


# ═══════════════════════════════════════════════════
# SLIDE 5: KI-Chatbot (erweitert)
# ═══════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)
slide_header(slide5, "KI-Integration", "Intelligenter Reiseberater")

# Left: Chat Flow
flows = [
    ("Natuerliche Sprache", "'Ich will nach Thailand' — der Bot erkennt Ziel, extrahiert IATA-Codes, startet Suche.", ACCENT_BLUE),
    ("Structured Output", "JSON Mode erzwingt valide API-Antworten fuer Flug-/Hotel-Datenextraktion.", ACCENT_GREEN),
    ("10+ Intent-Typen", "Flug, Hotel, Paket, Inspiration, Vergleich, Packliste, Budget, Visa, Wetter, Planung.", ACCENT_PURPLE),
    ("Suggested Replies", "Kontextbezogene Quick-Reply-Buttons nach jeder Bot-Antwort.", ACCENT_ORANGE),
    ("Auto Visa-Check", "Bei jeder Reiseanfrage automatischer Visa-Check + Auswaertiges Amt Link.", ACCENT_PINK),
]

y = 2.2
for i, (title, desc, color) in enumerate(flows):
    num_shape = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.05), Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    add_text(slide5, str(i + 1), Inches(0.8), Inches(y + 0.07), Inches(0.45), Inches(0.45),
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    card = add_shape(slide5, Inches(1.5), Inches(y), Inches(11.1), Inches(0.9), GLASS_BG, 0.04)
    add_text(slide5, title, Inches(1.8), Inches(y + 0.08), Inches(3), Inches(0.35),
             font_size=16, color=WHITE, bold=True)
    add_text(slide5, desc, Inches(1.8), Inches(y + 0.45), Inches(10.5), Inches(0.4),
             font_size=12, color=TEXT_DIM)
    y += 1.0


# ═══════════════════════════════════════════════════
# SLIDE 6: Prompt Engineering & Structured Output
# ═══════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6)
slide_header(slide6, "Prompt Engineering", "Structured Output & Intent-Erkennung")

# Left: System Prompt
add_text(slide6, "System-Prompt Design", Inches(0.8), Inches(2.1), Inches(5), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

prompt_features = [
    "Reise-Buddy Persona (locker, hilfsbereit)",
    "Proaktives Verhalten (wie Booking.com & Kayak)",
    "Paket-Empfehlungen (Flug + Hotel Kombination)",
    "Saisonale Tipps & Geheimtipps",
    "Kulturelle Hinweise fuer Ziellaender",
    "Kontext-Merken (Budget, Praeferenzen, Name)",
    "Max 2-3 Saetze (Token-effizient)",
]

y = 2.7
for item in prompt_features:
    row = add_shape(slide6, Inches(0.8), Inches(y), Inches(5.5), Inches(0.42), GLASS_BG, 0.04)
    add_text(slide6, "  " + item, Inches(1.0), Inches(y + 0.04), Inches(5.1), Inches(0.38),
             font_size=12, color=TEXT_DIM)
    dot = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.88), Inches(y + 0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE
    dot.line.fill.background()
    y += 0.48

# Right: Structured Output
add_text(slide6, "Structured Output (JSON Mode)", Inches(7), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

so_card = add_shape(slide6, Inches(7), Inches(2.7), Inches(5.5), Inches(2.4), GLASS_BG, 0.06)
add_text(slide6, "response_format = {\"type\": \"json_object\"}", Inches(7.3), Inches(2.85), Inches(5), Inches(0.4),
         font_size=13, color=ACCENT_GREEN, bold=True)
add_text(slide6, "Vorteile gegenueber Standard-Output:", Inches(7.3), Inches(3.3), Inches(5), Inches(0.3),
         font_size=12, color=WHITE, bold=True)

so_items = [
    "Garantiert valides JSON (kein Regex-Parsing noetig)",
    "95% Genauigkeit bei Datenextraktion (vs. 60% Zero-Shot)",
    "temperature=0.3 fuer konsistente JSON-Antworten",
    "temperature=0.7 fuer natuerlichen Chat",
]
add_bullet_list(slide6, so_items, Inches(7.3), Inches(3.65), Inches(5), Inches(1.3), font_size=11, color=TEXT_DIM)

# Comparison Box
add_text(slide6, "Modell-Vergleich", Inches(7), Inches(5.3), Inches(5.5), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

comp_card = add_shape(slide6, Inches(7), Inches(5.8), Inches(5.5), Inches(1.2), GLASS_BG, 0.06)
models = [
    ("GPT-4o-mini", "Gewaehlt", "$0.00015/1K", "0.5-1.5s", ACCENT_GREEN),
    ("GPT-4o", "Zu teuer", "$0.005/1K", "2-4s", TEXT_MUTED),
    ("GPT-3.5", "Schwach DE", "$0.0005/1K", "0.3-1s", TEXT_MUTED),
]

add_text(slide6, "Modell          Status          Kosten/1K       Speed", Inches(7.2), Inches(5.85), Inches(5.2), Inches(0.3),
         font_size=10, color=TEXT_MUTED, bold=True)

y_m = 6.15
for name, status, cost, speed, col in models:
    add_text(slide6, f"{name}       {status}        {cost}      {speed}", Inches(7.2), Inches(y_m), Inches(5.2), Inches(0.25),
             font_size=10, color=col)
    y_m += 0.25


# ═══════════════════════════════════════════════════
# SLIDE 7: Datenbank Schema
# ═══════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7)
slide_header(slide7, "Datenbank", "SQLite Schema (4 Tabellen)")

tables = [
    ("SearchHistory", ACCENT_BLUE, [
        ("id", "INTEGER", "PK, AUTO"),
        ("from_airport", "TEXT", "IATA-Code"),
        ("to_airport", "TEXT", "IATA-Code"),
        ("date", "TEXT", "YYYY-MM-DD"),
        ("timestamp", "DATETIME", "UTC Default"),
    ]),
    ("ChatCache", ACCENT_PURPLE, [
        ("id", "INTEGER", "PK, AUTO"),
        ("question", "TEXT", "INDEXED"),
        ("answer", "TEXT", "Gecachte Antwort"),
        ("created_at", "DATETIME", "UTC Default"),
    ]),
    ("CartItem", ACCENT_GREEN, [
        ("id", "INTEGER", "PK, AUTO"),
        ("session_id", "TEXT", "INDEXED"),
        ("item_type", "TEXT", "flight/hotel/activity"),
        ("item_data", "TEXT", "JSON String"),
        ("price", "REAL", "Preis in EUR"),
        ("currency", "TEXT", "Default EUR"),
        ("quantity", "INTEGER", "Default 1"),
        ("added_at", "DATETIME", "UTC Default"),
    ]),
    ("Flight", ACCENT_ORANGE, [
        ("id", "INTEGER", "PK, AUTO"),
        ("origin", "TEXT", "Abflughafen"),
        ("destination", "TEXT", "Zielflughafen"),
        ("departure_date", "TEXT", "Datum"),
        ("airline", "TEXT", "Nullable"),
        ("price", "REAL", "Nullable"),
        ("stops", "INTEGER", "Nullable"),
    ]),
]

col_x = 0.5
for tbl_name, tbl_color, fields in tables:
    w = 3.0
    h = 0.3 + len(fields) * 0.32 + 0.1
    card = add_shape(slide7, Inches(col_x), Inches(2.2), Inches(w), Inches(h), GLASS_BG, 0.04)
    # Header
    hdr = add_shape(slide7, Inches(col_x), Inches(2.2), Inches(w), Inches(0.45), tbl_color, 0.04)
    add_text(slide7, tbl_name, Inches(col_x), Inches(2.22), Inches(w), Inches(0.45),
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Fields
    fy = 2.75
    for fname, ftype, fdesc in fields:
        add_text(slide7, fname, Inches(col_x + 0.1), Inches(fy), Inches(1.1), Inches(0.28),
                 font_size=9, color=WHITE, bold=True)
        add_text(slide7, ftype, Inches(col_x + 1.2), Inches(fy), Inches(0.7), Inches(0.28),
                 font_size=9, color=tbl_color)
        add_text(slide7, fdesc, Inches(col_x + 1.9), Inches(fy), Inches(1.0), Inches(0.28),
                 font_size=8, color=TEXT_MUTED)
        fy += 0.32
    col_x += 3.15

# In-Memory note
add_text(slide7, "Zusaetzlich: In-Memory User-DB (Demo-Auth) + Visa-Datenbank (45+ Laender) + Airport-DB (50+ Flughaefen)",
         Inches(0.8), Inches(6.5), Inches(11), Inches(0.5), font_size=12, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════
# SLIDE 8: Security & Testing
# ═══════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8)
slide_header(slide8, "Security & Testing", "Produktionsreife Absicherung")

# Security (Left)
add_text(slide8, "Security-Massnahmen", Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=18, color=ACCENT_PINK, bold=True)

security_items = [
    ("CORS Restricted", "Nur localhost Origins erlaubt (kein allow_origins=*)"),
    ("Security Headers", "X-Frame-Options, X-XSS-Protection, CSP, Referrer-Policy"),
    ("Input Sanitization", "bleach.clean() gegen XSS, max_length Limits"),
    ("Passwort-Hashing", "bcrypt via passlib (kein Klartext)"),
    ("Rate Limiting", "10-30 req/min pro Endpoint (SlowAPI)"),
    ("Session Validation", "UUID-Format, Laengen-Check, Sonderzeichen-Filter"),
    ("Pydantic Validators", "Typ-Pruefung, Wertebereiche, Format-Validierung"),
    ("Env-Var Validation", "Startup-Check fuer alle API Keys"),
]

y = 2.7
for title, desc in security_items:
    row = add_shape(slide8, Inches(0.8), Inches(y), Inches(5.5), Inches(0.52), GLASS_BG, 0.04)
    add_text(slide8, title, Inches(1.0), Inches(y + 0.02), Inches(2.2), Inches(0.25),
             font_size=11, color=ACCENT_PINK, bold=True)
    add_text(slide8, desc, Inches(1.0), Inches(y + 0.26), Inches(5.1), Inches(0.22),
             font_size=9, color=TEXT_DIM)
    y += 0.57

# Testing (Right)
add_text(slide8, "Test Coverage", Inches(7), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=18, color=ACCENT_GREEN, bold=True)

test_files = [
    ("test_api.py", "23 KB", "API Endpoint Tests + Edge Cases", ACCENT_BLUE),
    ("test_services.py", "6 KB", "Visa, Airport, Warn-Level Tests", ACCENT_PURPLE),
    ("test_security.py", "17 KB", "XSS, CORS, SQL-Injection, Auth", ACCENT_PINK),
    ("test_integration.py", "17 KB", "Booking Flow, Chat History, Visa", ACCENT_ORANGE),
]

y = 2.7
for fname, size, desc, color in test_files:
    card = add_shape(slide8, Inches(7), Inches(y), Inches(5.5), Inches(0.95), GLASS_BG, 0.06)
    ind = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(y), Pt(4), Inches(0.95))
    ind.fill.solid()
    ind.fill.fore_color.rgb = color
    ind.line.fill.background()
    add_text(slide8, fname, Inches(7.3), Inches(y + 0.08), Inches(3), Inches(0.3),
             font_size=14, color=WHITE, bold=True)
    add_text(slide8, size, Inches(11.5), Inches(y + 0.08), Inches(0.8), Inches(0.3),
             font_size=11, color=color, bold=True)
    add_text(slide8, desc, Inches(7.3), Inches(y + 0.45), Inches(4.8), Inches(0.4),
             font_size=11, color=TEXT_DIM)
    y += 1.1

add_text(slide8, "Total: 63 KB Testcode  |  4 Test-Suites  |  50+ Test Cases",
         Inches(7), Inches(6.6), Inches(5.5), Inches(0.4),
         font_size=12, color=ACCENT_GREEN, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 9: Design System
# ═══════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9)
slide_header(slide9, "Design", "Apple Liquid Glass (iOS 26)")

# Color Palette
add_text(slide9, "Farbpalette", Inches(0.8), Inches(2.1), Inches(3), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

colors = [
    ("#5AC8FA", "Accent Blue", ACCENT_BLUE),
    ("#BF5AF2", "Accent Purple", ACCENT_PURPLE),
    ("#FF375F", "Accent Pink", ACCENT_PINK),
    ("#FF9F0A", "Accent Orange", ACCENT_ORANGE),
    ("#30D158", "Accent Green", ACCENT_GREEN),
    ("#0A0A0F", "Background", RGBColor(0x20, 0x20, 0x30)),
]

x = 0.8
for hex_code, name, color in colors:
    swatch = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.6), Inches(1.85), Inches(1.2))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = color
    swatch.line.fill.background()
    swatch.adjustments[0] = 0.1
    add_text(slide9, hex_code, Inches(x), Inches(3.85), Inches(1.85), Inches(0.3),
             font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide9, name, Inches(x), Inches(4.1), Inches(1.85), Inches(0.3),
             font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 2.0

# Design Principles
add_text(slide9, "Design-Prinzipien", Inches(0.8), Inches(4.7), Inches(3), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

principles = [
    ("Glasoberflaechen", "backdrop-filter: blur(40px)\ntransluzente Hintergruende", ACCENT_BLUE),
    ("Prismatische Raender", "Regenbogen-Lichtbrechung\nan Elementkanten", ACCENT_PURPLE),
    ("Ambient Orbs", "Grosse, weiche Farbkugeln\nim Hintergrund", ACCENT_PINK),
    ("Responsive + A11y", "Mobile/Tablet + ARIA\nLabels + Focus States", ACCENT_GREEN),
]

x = 0.8
for title, desc, color in principles:
    card = add_shape(slide9, Inches(x), Inches(5.2), Inches(2.95), Inches(1.6), GLASS_BG, 0.06)
    line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.2), Inches(2.95), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    add_text(slide9, title, Inches(x + 0.2), Inches(5.4), Inches(2.55), Inches(0.35),
             font_size=14, color=WHITE, bold=True)
    add_text(slide9, desc, Inches(x + 0.2), Inches(5.8), Inches(2.55), Inches(0.8),
             font_size=11, color=TEXT_DIM)
    x += 3.1


# ═══════════════════════════════════════════════════
# SLIDE 10: GenAI Ethics & Data Preprocessing
# ═══════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10)
slide_header(slide10, "GenAI Ethics & Data", "Verantwortungsvolle KI-Nutzung")

# Ethics (Left)
add_text(slide10, "Ethik-Massnahmen", Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=18, color=ACCENT_ORANGE, bold=True)

ethics = [
    ("Reisewarnungen", "Automatische Warnung bei kritischen Laendern (Syrien, Afghanistan, etc.)"),
    ("Rate Limiting", "Verhindert API-Missbrauch und schuetzt vor Ueberlastung"),
    ("Transparente Quellen", "Links zum Auswaertigen Amt statt eigene 'Fakten'"),
    ("Datenschutz", "Keine persistente User-Speicherung, Session-basiert"),
    ("Token-Effizienz", "Lokale Lookups statt KI-Calls wo moeglich (0 Tokens fuer Visa)"),
    ("Keine Halluzinationen", "Visa-Daten aus lokaler DB, nicht von GPT erfunden"),
]

y = 2.7
for title, desc in ethics:
    row = add_shape(slide10, Inches(0.8), Inches(y), Inches(5.5), Inches(0.65), GLASS_BG, 0.04)
    add_text(slide10, title, Inches(1.0), Inches(y + 0.04), Inches(2.5), Inches(0.28),
             font_size=12, color=ACCENT_ORANGE, bold=True)
    add_text(slide10, desc, Inches(1.0), Inches(y + 0.32), Inches(5.1), Inches(0.28),
             font_size=10, color=TEXT_DIM)
    y += 0.72

# Data Preprocessing (Right)
add_text(slide10, "Data Preprocessing", Inches(7), Inches(2.1), Inches(5.5), Inches(0.4),
         font_size=18, color=ACCENT_BLUE, bold=True)

preproc = [
    ("Input Sanitization", "bleach.clean() entfernt HTML-Tags und XSS-Vektoren"),
    ("Session Validation", "Regex-Filter: nur alphanumerisch + Bindestrich/Unterstrich"),
    ("Pydantic Validators", "Typ-Pruefung, Wertebereiche (Preis 0-1M, Rolle in user/assistant)"),
    ("Date Parsing", "datetime.strptime() validiert YYYY-MM-DD Format"),
    ("Country Normalization", "Aliase aufloesen: 'USA'/'Amerika'/'United States' -> 'usa'"),
    ("Chat History Limit", "Max 50 Nachrichten, max 5000 Zeichen pro Nachricht"),
]

y = 2.7
for title, desc in preproc:
    row = add_shape(slide10, Inches(7), Inches(y), Inches(5.5), Inches(0.65), GLASS_BG, 0.04)
    add_text(slide10, title, Inches(7.2), Inches(y + 0.04), Inches(2.5), Inches(0.28),
             font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide10, desc, Inches(7.2), Inches(y + 0.32), Inches(5.1), Inches(0.28),
             font_size=10, color=TEXT_DIM)
    y += 0.72


# ═══════════════════════════════════════════════════
# SLIDE 11: API Endpunkte
# ═══════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11)
slide_header(slide11, "Backend", "20+ REST API Endpunkte")

header_bg = add_shape(slide11, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.5), ACCENT_BLUE, 0.02)
add_text(slide11, "Methode", Inches(1), Inches(2.12), Inches(1.2), Inches(0.5),
         font_size=12, color=WHITE, bold=True)
add_text(slide11, "Endpunkt", Inches(2.3), Inches(2.12), Inches(4), Inches(0.5),
         font_size=12, color=WHITE, bold=True)
add_text(slide11, "Beschreibung", Inches(6.5), Inches(2.12), Inches(4), Inches(0.5),
         font_size=12, color=WHITE, bold=True)
add_text(slide11, "Rate Limit", Inches(10.5), Inches(2.12), Inches(1.8), Inches(0.5),
         font_size=12, color=WHITE, bold=True)

endpoints = [
    ("GET", "/health", "Health Check", "-"),
    ("GET", "/airports?q=", "Flughafen-Autocomplete", "30/min"),
    ("GET", "/search_flights", "Flugsuche via Skyscanner", "10/min"),
    ("GET", "/search_hotels", "Hotelsuche", "10/min"),
    ("GET", "/search_destination", "Hotel-Ziel Autocomplete", "30/min"),
    ("POST", "/chat", "KI-Chat mit History + Intents", "20/min"),
    ("GET", "/smart_ask", "Einfacher Chat (Legacy)", "10/min"),
    ("POST", "/cart/add", "Warenkorb: Hinzufuegen", "30/min"),
    ("GET", "/cart/{session_id}", "Warenkorb: Abrufen", "30/min"),
    ("DEL", "/cart/remove", "Warenkorb: Entfernen", "30/min"),
    ("DEL", "/cart/clear/{id}", "Warenkorb: Leeren", "10/min"),
    ("POST", "/auth/login", "Benutzer-Login (bcrypt)", "-"),
    ("POST", "/auth/register", "Registrierung", "-"),
    ("GET", "/visa/{country}", "Visa-Informationen", "30/min"),
    ("GET", "/travel-info/{country}", "Kombinierte Reiseinfo", "10/min"),
]

y = 2.7
for i, (method, endpoint, desc, rate) in enumerate(endpoints):
    bg_color = GLASS_BG if i % 2 == 0 else RGBColor(0x14, 0x14, 0x20)
    add_shape(slide11, Inches(0.8), Inches(y), Inches(11.7), Inches(0.32), bg_color, 0.01)
    method_color = ACCENT_GREEN if method == "GET" else ACCENT_ORANGE if method == "POST" else ACCENT_PINK
    add_text(slide11, method, Inches(1), Inches(y + 0.01), Inches(1.2), Inches(0.3),
             font_size=10, color=method_color, bold=True)
    add_text(slide11, endpoint, Inches(2.3), Inches(y + 0.01), Inches(4), Inches(0.3),
             font_size=10, color=WHITE)
    add_text(slide11, desc, Inches(6.5), Inches(y + 0.01), Inches(4), Inches(0.3),
             font_size=10, color=TEXT_DIM)
    add_text(slide11, rate, Inches(10.5), Inches(y + 0.01), Inches(1.8), Inches(0.3),
             font_size=10, color=TEXT_MUTED)
    y += 0.33


# ═══════════════════════════════════════════════════
# SLIDE 12: Roadmap
# ═══════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12)
slide_header(slide12, "Zukunft", "Roadmap & Milestones")

done_items = [
    "Flugsuche mit Skyscanner API",
    "KI-Chatbot mit GPT-4o-mini + Structured Output",
    "Hotelsuche mit Sky-Scrapper API",
    "Visa-Informationen + Reisewarnungen (45+ Laender)",
    "Warenkorb-System (Fluege, Hotels, Aktivitaeten)",
    "Benutzer-Authentifizierung (bcrypt)",
    "Apple Liquid Glass Design System",
    "Mietwagen-Seite (funktional)",
    "Aktivitaeten-Seite (funktional)",
    "Security: CORS, Headers, XSS-Schutz, Sanitization",
    "50+ Unit/Integration/Security Tests",
]

add_text(slide12, "Erledigt (V1)", Inches(0.8), Inches(2.2), Inches(3), Inches(0.4),
         font_size=16, color=ACCENT_GREEN, bold=True)

y = 2.7
for item in done_items:
    row = add_shape(slide12, Inches(0.8), Inches(y), Inches(5.5), Inches(0.37), GLASS_BG, 0.04)
    add_text(slide12, "  " + item, Inches(1.05), Inches(y + 0.02), Inches(5), Inches(0.35),
             font_size=10, color=TEXT_DIM)
    check = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.88), Inches(y + 0.06), Inches(0.24), Inches(0.24))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_GREEN
    check.line.fill.background()
    y += 0.39

planned_items = [
    "Zahlungsintegration (Stripe)",
    "Echte Mietwagen-API Anbindung",
    "Push-Benachrichtigungen (Preisalarm)",
    "Progressive Web App (PWA)",
    "Mehrsprachigkeit (EN, FR, ES)",
    "User-Profile & Buchungsverlauf",
]

add_text(slide12, "Geplant (V2)", Inches(7), Inches(2.2), Inches(3), Inches(0.4),
         font_size=16, color=ACCENT_ORANGE, bold=True)

y = 2.7
for item in planned_items:
    row = add_shape(slide12, Inches(7), Inches(y), Inches(5.5), Inches(0.37), GLASS_BG, 0.04)
    add_text(slide12, "  " + item, Inches(7.25), Inches(y + 0.02), Inches(5), Inches(0.35),
             font_size=10, color=TEXT_DIM)
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.08), Inches(y + 0.06), Inches(0.24), Inches(0.24))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GLASS_BG
    circle.line.color.rgb = ACCENT_ORANGE
    circle.line.width = Pt(2)
    y += 0.39

# Stats Box
add_text(slide12, "Projekt-Statistiken", Inches(7), Inches(5.2), Inches(5.5), Inches(0.4),
         font_size=16, color=WHITE, bold=True)

stats = [
    ("1.345", "Zeilen Backend (main.py)"),
    ("4.495", "Zeilen CSS (style.css)"),
    ("2.185", "Zeilen JavaScript (app.js)"),
    ("20+", "REST API Endpunkte"),
    ("4", "Test-Suites / 50+ Tests"),
    ("45+", "Laender mit Visa-Daten"),
]

sx = 7.0
sy = 5.7
for val, label in stats:
    card = add_shape(slide12, Inches(sx), Inches(sy), Inches(1.75), Inches(0.75), GLASS_BG, 0.08)
    add_text(slide12, val, Inches(sx), Inches(sy + 0.05), Inches(1.75), Inches(0.35),
             font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide12, label, Inches(sx), Inches(sy + 0.4), Inches(1.75), Inches(0.3),
             font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    sx += 1.85
    if sx > 11.5:
        sx = 7.0
        sy += 0.85


# ═══════════════════════════════════════════════════
# SLIDE 13: Abschluss
# ═══════════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13)

c1 = slide13.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(5), Inches(5))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x2E)
c1.line.fill.background()

c2 = slide13.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(4), Inches(4))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(0x1A, 0x0F, 0x28)
c2.line.fill.background()

add_text(slide13, "Vielen Dank!", Inches(1), Inches(2), Inches(11.3), Inches(1.2),
         font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide13, Inches(5.6), Inches(3.3), Inches(2.1), ACCENT_BLUE)

add_text(slide13, "Fragen & Live Demo", Inches(1), Inches(3.7), Inches(11.3), Inches(0.7),
         font_size=26, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

info = [
    ("GitHub", "github.com/sakhiaryan/FlyAi", ACCENT_BLUE),
    ("Live Demo", "127.0.0.1:5500", ACCENT_GREEN),
    ("API Docs", "127.0.0.1:8000/docs", ACCENT_PURPLE),
]

x = 2.5
for title, link, color in info:
    card = add_shape(slide13, Inches(x), Inches(4.8), Inches(2.7), Inches(0.9), GLASS_BG, 0.08)
    line = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.8), Inches(2.7), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    add_text(slide13, title, Inches(x), Inches(4.9), Inches(2.7), Inches(0.35),
             font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide13, link, Inches(x), Inches(5.25), Inches(2.7), Inches(0.3),
             font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 2.9

add_text(slide13, "Aryan Sakhi  |  Masterschool GenAI  |  2026", Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SPEICHERN
# ═══════════════════════════════════════════════════
output_path = "/Users/aryan/Desktop/FlyAi/FlyAI_Praesentation.pptx"
prs.save(output_path)
print(f"Praesentation gespeichert: {output_path}")
print(f"Slides: {len(prs.slides)}")
